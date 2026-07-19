#!/usr/bin/env python3
"""Generate the CLSSGB Lesson Plan (LP) DOCX in the Tertiary house format.

Cover page + Document Version Control Record + auto TOC + Arial 11pt body +
colour-coded 4-day schedule tables (9:30am-6:30pm, 8 training hours/day, 1h
lunch, tea within, final assessment Day 4). Topics/labs come from
course_data + the domain data files so the LP stays aligned with the deck,
guide and labs.
"""
import os, sys
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

HERE=os.path.dirname(os.path.abspath(__file__)); sys.path.insert(0,HERE)
import course_data as C
from data_domain1 import DOMAIN1; from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3; from data_domain4 import DOMAIN4
from data_domain5 import DOMAIN5
ACT=DOMAIN1+DOMAIN2+DOMAIN3+DOMAIN4+DOMAIN5
import prodoc
def _find_repo(start):
    env=os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env): return env
    d=start
    for _ in range(8):
        d=os.path.dirname(d)
        if os.path.isdir(os.path.join(d,"courseware")) and os.path.isdir(os.path.join(d,"labs")): return d
    return os.path.dirname(os.path.dirname(HERE))
REPO=_find_repo(HERE); ASSETS=os.path.join(os.path.dirname(HERE),"assets")

BRAND=RGBColor(0x1F,0x6F,0xEB); DARK=RGBColor(0x11,0x18,0x27); GREY=RGBColor(0x55,0x5B,0x66)
HEADER_FILL="1F6FEB"; TOPIC_FILL="E8F0FE"; BREAK_FILL="FFF4E5"; LUNCH_FILL="FDE9D9"; ASSESS_FILL="E8F7EE"

def lab_titles(nums):
    return "; ".join(f"Lab {a['num']}: {a['title']}" for a in ACT if a['num'] in nums)

# ------------------------------------------------ slide ranges (read from the built deck)
def _scan_deck():
    """Build ONE ordered index of every section start and every lab start in the
    deck. Ranges are then derived by taking the next boundary of ANY kind, so a
    section's teaching range stops where its labs begin, and a lab range stops
    where the next section begins."""
    try:
        import re
        from pptx import Presentation
        deck = os.path.join(REPO, "courseware", f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
        prs = Presentation(deck)
    except Exception:
        return [], 0
    MARKERS = {"COURSE ADMINISTRATION": "Admin", "FOUNDATIONS": "Foundations",
               "DMAIC · DEFINE": "Define", "DMAIC · MEASURE": "Measure",
               "DMAIC · ANALYZE": "Analyze", "DMAIC · IMPROVE": "Improve",
               "DMAIC · CONTROL": "Control", "WRAP-UP": "Wrap-up"}
    marks = []
    for i, s in enumerate(prs.slides, 1):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if t in MARKERS:
                marks.append((i, "section", MARKERS[t])); break
            m = re.fullmatch(r"LAB (\d+)", t)
            if m:
                num = int(m.group(1))
                if not any(k == "lab" and v == num for _, k, v in marks):
                    marks.append((i, "lab", num))
                break
    marks.sort(key=lambda x: x[0])
    return marks, len(prs.slides._sldIdLst)


MARKS, DECK_TOTAL = _scan_deck()
BOUNDS = [m[0] for m in MARKS]


def _end_after(start):
    """Last slide before the next boundary of any kind."""
    later = [b for b in BOUNDS if b > start]
    return (min(later) - 1) if later else DECK_TOTAL


def _slide_ranges():
    """Teaching range for each section: from its section slide up to (but not
    including) whatever comes next — its first lab, or the next section."""
    out = {}
    for start, kind, name in MARKS:
        if kind == "section":
            out[name] = f"{start}–{_end_after(start)}"
    return out

SLIDES = _slide_ranges()


LAB_STARTS = {v: s for s, k, v in MARKS if k == "lab"}


def sll(nums):
    """Slide range spanning the given labs. The range ends at the slide before
    the next boundary of ANY kind, so it never spills into the following
    section (e.g. the closing/assessment block)."""
    starts = [LAB_STARTS[n] for n in nums if n in LAB_STARTS]
    if not starts:
        return ""
    return f"  [slides {min(starts)}–{_end_after(max(starts))}]"

def sl(*names):
    """Slide-range annotation for the schedule, e.g. '  [slides 46-102]'.

    Accepts several section names and merges them into ONE contiguous range, so
    a row covering two sections never shows two bracketed ranges."""
    parts = [SLIDES[n] for n in names if n in SLIDES]
    if not parts:
        return ""
    if len(parts) == 1:
        return f"  [slides {parts[0]}]"
    lo = parts[0].split("–")[0]
    hi = parts[-1].split("–")[-1]
    return f"  [slides {lo}–{hi}]"

# ------------------------------------------------ schedule (single source of truth for timing)
# (start, end, minutes, kind, activity_text)  kind: admin/topic/lab/break/lunch/assess/recap
SCHEDULE = {
 1: (C.DAY_THEMES[1], [
    ("9:30","10:00",30,"admin","Welcome, course introduction, ground rules and mandatory digital attendance (AM)"+sl("Admin")),
    ("10:00","11:15",75,"topic","FOUNDATIONS — What is Quality; Lean vs Six Sigma vs Lean Six Sigma; the five Lean principles; history of Six Sigma; Y = f(X); classical vs Six Sigma quality"+sl("Foundations")),
    ("11:15","11:30",15,"break","Tea break"),
    ("11:30","12:30",60,"topic","FOUNDATIONS — sigma level and DPMO; the 1.5 sigma shift; Cost of Poor Quality; belt roles and what a Green Belt leads; project selection and the cause/solution matrix; the DMAIC roadmap and tollgates"),
    ("12:30","13:30",60,"lunch","Lunch break"),
    ("13:30","15:00",90,"lab","Hands-on: "+lab_titles([1,2])+sll([1,2])),
    ("15:00","15:15",15,"break","Tea break"),
    ("15:15","16:45",90,"topic","DMAIC · DEFINE — Voice of the Customer; affinity diagrams; Kano analysis; CTQ trees and operational definitions; the project charter; problem and goal statements; scope; SIPOC; stakeholders and RACI"+sl("Define")),
    ("16:45","18:00",75,"lab","Hands-on: "+lab_titles([3,4])+sll([3,4])),
    ("18:00","18:30",30,"recap","Day 1 recap, Q&A and PM digital attendance"),
 ]),
 2: (C.DAY_THEMES[2], [
    ("9:30","9:45",15,"recap","Day 1 recap and mandatory digital attendance (AM)"),
    ("9:45","11:00",75,"lab","Hands-on: "+lab_titles([5,6])+sll([5,6])),
    ("11:00","11:15",15,"break","Tea break"),
    ("11:15","12:30",75,"topic","DMAIC · MEASURE — process mapping symbols and swimlanes; handoff analysis; the eight wastes (DOWNTIME); value-added analysis; value stream mapping; lean metrics and takt time"+sl("Measure")),
    ("12:30","13:30",60,"lunch","Lunch break"),
    ("13:30","15:00",90,"lab","Hands-on: "+lab_titles([7,8])+sll([7,8])),
    ("15:00","15:15",15,"break","Tea break"),
    ("15:15","16:30",75,"topic","DMAIC · MEASURE — types of data; data collection plans; sampling techniques; sample size calculation; Measurement System Analysis and Gage R&R; yield, DPU, DPO, DPMO and RTY; descriptive statistics and baseline capability"),
    ("16:30","18:15",105,"lab","Hands-on: "+lab_titles([9,10,11])+sll([9,10,11])),
    ("18:15","18:30",15,"recap","Day 2 recap and PM digital attendance"),
 ]),
 3: (C.DAY_THEMES[3], [
    ("9:30","9:45",15,"recap","Day 2 recap and mandatory digital attendance (AM)"),
    ("9:45","11:00",75,"lab","Hands-on: "+lab_titles([12,13])+sll([12,13])),
    ("11:00","11:15",15,"break","Tea break"),
    ("11:15","12:30",75,"topic","DMAIC · ANALYZE — common vs special cause variation; tampering; run charts and the six non-random patterns; the Pareto principle; stratification; boxplots"+sl("Analyze")),
    ("12:30","13:30",60,"lunch","Lunch break"),
    ("13:30","14:45",75,"lab","Hands-on: "+lab_titles([14,15])+sll([14,15])),
    ("14:45","15:00",15,"break","Tea break"),
    ("15:00","16:30",90,"topic","DMAIC · ANALYZE — Fishbone (5M+E); 5 Whys; multi-voting; hypothesis testing (H0/Ha, alpha, Type I and II errors, test selection, p-value decision rule); correlation and the coefficient of determination; regression"),
    ("16:30","18:15",105,"lab","Hands-on: "+lab_titles([16,17,18])+sll([16,17,18])),
    ("18:15","18:30",15,"recap","Day 3 recap and PM digital attendance"),
 ]),
 4: (C.DAY_THEMES[4], [
    ("9:30","9:45",15,"recap","Day 3 recap and mandatory digital attendance (AM)"),
    ("9:45","10:45",60,"topic","DMAIC · IMPROVE — solution generation and brainwriting; benchmarking; the solution selection matrix; 5S; poka-yoke; pull, JIT and Heijunka; standard work; FMEA and RPN; Design of Experiments; piloting"+sl("Improve")),
    ("10:45","11:00",15,"break","Tea break"),
    ("11:00","12:30",90,"lab","Hands-on: "+lab_titles([19,20,21,22])+sll([19,20,21,22])),
    ("12:30","13:30",60,"lunch","Lunch break"),
    ("13:30","14:30",60,"topic","DMAIC · CONTROL — process capability Cp and Cpk; control limits vs specification limits; SPC and control chart selection; the eight out-of-control rules; the control plan; SOPs; visual management; verifying the gain and handover"+sl("Control")),
    ("14:30","15:15",45,"lab","Hands-on: "+lab_titles([23])+sll([23])),
    ("15:15","15:30",15,"break","Tea break"),
    ("15:30","15:45",15,"lab","Hands-on: "+lab_titles([24,25])+sll([24,25])),
    ("15:45","16:00",15,"assess","Course recap and Briefing for Assessment"+sl("Wrap-up")),
    ("16:00","17:00",60,"assess","Written Assessment (WA) — Short-Answer Questions (SAQ), 2 questions, 60 minutes, open book"),
    ("17:00","18:30",90,"assess","Practical Performance (PP) — 3 applied DMAIC tasks, 90 minutes, open book. PM digital attendance"),
 ]),
}

# ------------------------------------------------ build document
doc=Document()
normal=doc.styles["Normal"]; normal.font.name="Arial"; normal.font.size=Pt(11)
prodoc.style_headings(doc)

prodoc.add_cover_page(doc,"LESSON PLAN",C.TITLE,C.VERSION.lstrip("v"),
                      org_logo=os.path.join(ASSETS,"tertiary-infotech-logo.png"),
                      course_logo=None, course_code=C.COURSE_CODE)
prodoc.add_version_control(doc,[
 ("1","1 June 2026","Initial release - CLSSGB lesson plan.",C.TRAINER),
 ("2","19 July 2026","Major revision: rebuilt as a 4-day Green Belt progression from the Yellow Belt "
  "(CLSSYB). Content deepened to the CSSC Green Belt scope (Chapters 1-24) - adds "
  "sampling and sample size calculation, Measurement System Analysis and Gage R&R, hypothesis testing "
  "with p-values and Type I/II errors, correlation and regression, FMEA and RPN, Design of Experiments, "
  "SPC control chart selection with the eight out-of-control rules, and process capability Cp/Cpk. "
  "Expanded from 10 to 25 hands-on labs across the five DMAIC phases; teaching diagrams imported from "
  "the original v21 trainer deck; schedule and labs realigned to the v6 WA (SAQ) and PP assessment papers.",C.TRAINER),
 ("3","19 July 2026","Removed the Practice Exam slide and its asset - there is no Six Sigma practice exam on exams.tertiaryinfotech.com.",C.TRAINER),
 ("4",C.VERSION_DATE,"Added the SIPOC & Process Map Builder (alfredang.github.io/sipoc) to the toolkit - wired into Lab 6 (SIPOC build, pain points, 'Check my SIPOC' validation) and Lab 7 (swimlane and handoff table), the Define phase slides, and a new Your Interactive Toolkit slide that introduces all five browser tools.",C.TRAINER),
])
prodoc.add_toc(doc)

def H(text,level=1):
    h=doc.add_heading(text,level=level); return h

H("Course Information",1)
info=[("Course Title",C.TITLE),("WSQ Course Reference",C.COURSE_CODE),
      ("Training Provider",C.ORG+"  ("+C.UEN.replace('UEN: ','UEN ')+")"),
      ("Duration","4 days · 8 training hours per day (32 hours)"),
      ("Daily Timing","9:30 am – 6:30 pm (1-hour lunch; tea breaks within training time)"),
      ("Mode","Instructor-led, hands-on Lean Six Sigma labs using the Northwind Retail Distribution Centre improvement scenario"),
      ("TSC Alignment",f"{C.TSC_TITLE} ({C.TSC_CODE})"),
      ("Trainer",C.TRAINER)]
t=doc.add_table(rows=0,cols=2); t.style="Table Grid"
for k,v in info:
    c=t.add_row().cells; c[0].text=""; r=c[0].paragraphs[0].add_run(k); r.bold=True; r.font.size=Pt(10)
    prodoc._shade_cell(c[0],TOPIC_FILL)
    c[1].text=""; c[1].paragraphs[0].add_run(v).font.size=Pt(10)

H("Learning Outcomes",1)
doc.add_paragraph("On completion of this course, learners will be able to:")
for lo in C.LEARNING_OUTCOMES:
    p=doc.add_paragraph(style="List Bullet"); p.add_run(lo).font.size=Pt(10.5)

H("Assessment",1)
for a in [C.ASSESSMENT["written"],C.ASSESSMENT["practical"],
          "Format: Open Book — course slides, Learner Guide and approved materials only.",
          "Final assessment is conducted on Day 4 from 4:00 pm.",C.ASSESSMENT["note"]]:
    p=doc.add_paragraph(style="List Bullet"); p.add_run(a).font.size=Pt(10.5)

def set_cell(cell,text,bold=False,size=9.5,color=None,fill=None,align=None):
    cell.text=""; p=cell.paragraphs[0]
    if align: p.alignment=align
    r=p.add_run(text); r.bold=bold; r.font.size=Pt(size); r.font.name="Arial"
    if color: r.font.color.rgb=color
    if fill: prodoc._shade_cell(cell,fill)

KIND_FILL={"topic":TOPIC_FILL,"break":BREAK_FILL,"lunch":LUNCH_FILL,"assess":ASSESS_FILL,
           "admin":"F3F5F8","recap":"F3F5F8","lab":None}

H("Course Schedule",1)
for day,(theme,rows) in SCHEDULE.items():
    H(f"Day {day} — {theme}",2)
    tbl=doc.add_table(rows=0,cols=3); tbl.style="Table Grid"; tbl.alignment=WD_TABLE_ALIGNMENT.CENTER
    hdr=tbl.add_row().cells
    for i,htext in enumerate(["Time","Duration","Topic / Activity"]):
        set_cell(hdr[i],htext,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
    training=0
    for start,end,mins,kind,text in rows:
        cells=tbl.add_row().cells; fill=KIND_FILL.get(kind)
        set_cell(cells[0],f"{start}–{end}",bold=(kind in ("topic","assess")),size=9.5,fill=fill)
        set_cell(cells[1],f"{mins} min",size=9.5,fill=fill)
        set_cell(cells[2],text,bold=(kind in ("topic","assess")),size=9.5,fill=fill)
        if kind!="lunch": training+=mins
    # widths
    for row in tbl.rows:
        row.cells[0].width=Inches(1.15); row.cells[1].width=Inches(0.9); row.cells[2].width=Inches(4.75)
    p=doc.add_paragraph(); r=p.add_run(f"Total training time: {training} minutes ({training//60} hours)."); r.italic=True; r.font.size=Pt(9.5); r.font.color.rgb=GREY
    assert training==480, f"Day {day} training minutes = {training}, expected 480"

H("Lab Reference (aligned to the DMAIC phases)",1)
tt=doc.add_table(rows=0,cols=3); tt.style="Table Grid"
hdr=tt.add_row().cells
for i,htext in enumerate(["DMAIC phase / Topic","Weighting","Labs"]):
    set_cell(hdr[i],htext,bold=True,size=10,color=RGBColor(0xFF,0xFF,0xFF),fill=HEADER_FILL)
for tp in C.TOPICS:
    acts=[a for a in ACT if a["topic"]==tp["num"]]
    cells=tt.add_row().cells
    set_cell(cells[0],f"{tp['phase']}: {tp['title']}",bold=True,size=9.5,fill=TOPIC_FILL)
    set_cell(cells[1],tp["weighting"],size=9.5,fill=TOPIC_FILL)
    set_cell(cells[2],", ".join(
        f"Lab {a['num']}" + (" (elective)" if a.get("elective") else "") for a in acts),size=9.5)

prodoc.add_page_numbers(doc)
prodoc.enable_update_fields(doc)
OUT=os.path.join(REPO,"courseware",f"LP-{C.SHORT_TITLE}.docx")
doc.save(OUT)
print("Saved",OUT)
