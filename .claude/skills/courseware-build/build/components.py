#!/usr/bin/env python3
"""Shared visual components for the CLSSYB deck (all-white Tertiary house style).

Base components (cover, section, content, two_col, cards3, big_statement,
tile_grid, flow_h, trainer_slide, brk) are the house set used across Tertiary WSQ
decks. Added here for Lean Six Sigma: concept diagrams drawn natively with
python-pptx so every key concept is explained VISUALLY rather than as bullets —
dmaic_wheel, sipoc_diagram, fishbone, pareto_chart, run_chart, normal_curve,
waste_wheel, compare_panels, vs_diagram, matrix2x2, ladder, timeline, formula_card.
"""
import math
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import copy

# ---------------- palette ----------------
BLUE   = RGBColor(0x1F, 0x6F, 0xEB)
TEAL   = RGBColor(0x10, 0xB9, 0x81)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xDC, 0x2C, 0x2C)
VIOLET = RGBColor(0x7C, 0x3A, 0xED)
INK    = RGBColor(0x16, 0x1B, 0x26)
GREY   = RGBColor(0x5B, 0x63, 0x72)
LIGHT  = RGBColor(0xF5, 0xF8, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LINE   = RGBColor(0xE2, 0xE8, 0xF0)
PALETTE = [BLUE, TEAL, VIOLET, AMBER]
# DMAIC phase colours, used consistently everywhere DMAIC appears
DMAIC_COLORS = [BLUE, TEAL, VIOLET, AMBER, RED]


def _default_repo_assets():
    """Locate <repo>/courseware/assets by walking up from this file.

    Mirrors the repo discovery used by the build_*.py entry points so a Deck
    built directly (e.g. in a test) still resolves its diagram assets.
    """
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return os.path.join(env, "courseware", "assets")
    d = os.path.dirname(os.path.abspath(__file__))
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return os.path.join(d, "courseware", "assets")
    return ""


class Deck:
    """Holds the presentation + page counter so components stay stateless."""

    def __init__(self, course, repo_assets=None):
        self.C = course
        # Absolute path to <repo>/courseware/assets — concept modules resolve the
        # imported trainer-deck diagrams (assets/diagrams/*) through this.
        self.REPO_ASSETS = repo_assets or _default_repo_assets()
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self.SW = self.prs.slide_width
        self.SH = self.prs.slide_height
        self.BLANK = self.prs.slide_layouts[6]
        self.page = 0

    # ---------------- primitives ----------------
    def slide(self):
        return self.prs.slides.add_slide(self.BLANK)

    def rect(self, s, x, y, w, h, color, line=None):
        sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, int(x), int(y), int(w), int(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line; sp.line.width = Pt(1)
        sp.shadow.inherit = False
        if sp.has_text_frame:
            sp.text_frame.text = ""
        return sp

    def roundrect(self, s, x, y, w, h, color, line=None):
        sp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y), int(w), int(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line; sp.line.width = Pt(1.25)
        sp.shadow.inherit = False
        if sp.has_text_frame:
            sp.text_frame.text = ""
        return sp

    def oval(self, s, x, y, w, h, color, line=None):
        sp = s.shapes.add_shape(MSO_SHAPE.OVAL, int(x), int(y), int(w), int(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line; sp.line.width = Pt(1.25)
        sp.shadow.inherit = False
        if sp.has_text_frame:
            sp.text_frame.text = ""
        return sp

    def chevron(self, s, x, y, w, h, color):
        sp = s.shapes.add_shape(MSO_SHAPE.CHEVRON, int(x), int(y), int(w), int(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = color
        sp.line.fill.background(); sp.shadow.inherit = False
        if sp.has_text_frame:
            sp.text_frame.text = ""
        return sp

    def line_h(self, s, x, y, w, color, width=1.5):
        c = s.shapes.add_connector(1, int(x), int(y), int(x + w), int(y))
        c.line.color.rgb = color; c.line.width = Pt(width)
        return c

    def line_seg(self, s, x1, y1, x2, y2, color, width=1.5, dash=False):
        c = s.shapes.add_connector(1, int(x1), int(y1), int(x2), int(y2))
        c.line.color.rgb = color; c.line.width = Pt(width)
        if dash:
            c.line._get_or_add_ln().append(
                _el('a:prstDash', {'val': 'dash'}))
        return c

    def txt(self, s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=4):
        """Render lines of runs into a textbox.

        A literal "\\n" inside a run becomes a vertical-tab line break in OOXML,
        which does NOT inherit paragraph alignment — continuation lines render
        left-shifted inside centred text. So expand any embedded newline into a
        real paragraph, which does carry the alignment.
        """
        expanded = []
        for ln in runs:
            parts = [[]]
            for t, sz, col, bold in ln:
                pieces = str(t).split("\n")
                for k, piece in enumerate(pieces):
                    if k:
                        parts.append([])
                    if piece:
                        parts[-1].append((piece, sz, col, bold))
            for part in parts:
                expanded.append(part)

        tb = s.shapes.add_textbox(int(x), int(y), int(w), int(h))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        for i, ln in enumerate(expanded):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.space_after = Pt(space)
            for t, sz, col, bold in ln:
                r = p.add_run(); r.text = t
                r.font.size = Pt(sz); r.font.bold = bold
                r.font.color.rgb = col; r.font.name = "Arial"
        return tb

    def bullets(self, s, x, y, w, h, items, size=18, color=INK, gap=10):
        tb = s.shapes.add_textbox(int(x), int(y), int(w), int(h))
        tf = tb.text_frame; tf.word_wrap = True
        for i, it in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(gap)
            lvl = it[1] if isinstance(it, tuple) else 0
            text = it[0] if isinstance(it, tuple) else it
            r = p.add_run(); r.text = ("•  " if lvl == 0 else "–  ") + text
            r.font.size = Pt(size if lvl == 0 else size - 2)
            r.font.color.rgb = color if lvl == 0 else GREY
            r.font.name = "Arial"
        return tb

    # ---------------- chrome ----------------
    def footer(self, s):
        self.page += 1
        C = self.C
        self.txt(s, Inches(0.4), Inches(7.05), Inches(7.5), Inches(0.35),
                 [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}", 9, GREY, False)]])
        self.txt(s, Inches(5.0), Inches(7.05), Inches(3.3), Inches(0.35),
                 [[("© 2026 Tertiary Infotech Academy Pte Ltd", 9, GREY, False)]], align=PP_ALIGN.CENTER)
        self.txt(s, Inches(12.4), Inches(7.05), Inches(0.6), Inches(0.35),
                 [[(str(self.page), 9, GREY, False)]], align=PP_ALIGN.RIGHT)

    def head(self, s, title, kicker=None, kcolor=BLUE):
        self.rect(s, 0, 0, self.SW, self.SH, WHITE)
        self.rect(s, 0, 0, Inches(0.28), Inches(1.55), kcolor)
        if kicker:
            self.txt(s, Inches(0.85), Inches(0.5), Inches(11.6), Inches(0.4), [[(kicker, 14, kcolor, True)]])
        self.txt(s, Inches(0.85), Inches(0.92), Inches(11.9), Inches(0.75), [[(title, 30, INK, True)]])
        return s

    # ---------------- base templates ----------------
    def cover(self, logo=None):
        C = self.C
        s = self.slide()
        self.rect(s, 0, 0, self.SW, self.SH, WHITE)
        self.rect(s, 0, 0, Inches(0.42), self.SH, BLUE)
        self.rect(s, Inches(0.9), Inches(1.55), Inches(2.4), Inches(0.14), TEAL)
        self.txt(s, Inches(0.9), Inches(1.95), Inches(11.6), Inches(2.3),
                 [[(C.TITLE, 42, INK, True)]])
        self.rect(s, Inches(0.92), Inches(4.35), Inches(2.4), Inches(0.06), TEAL)
        self.txt(s, Inches(0.9), Inches(4.65), Inches(12), Inches(1.4),
                 [[(f"WSQ Course Code: {C.COURSE_CODE}", 16, GREY, False)],
                  [(f"Conducted by {C.ORG}  ·  {C.UEN}", 14, GREY, False)],
                  [(f"TSC: {C.TSC_TITLE}  ·  {C.TSC_CODE}", 13, GREY, False)]], space=6)
        self.txt(s, Inches(0.9), Inches(6.42), Inches(12), Inches(0.4),
                 [[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}  ·  Trainer: {C.TRAINER}", 12, GREY, False)]])
        self.txt(s, Inches(0.9), Inches(6.85), Inches(12), Inches(0.34),
                 [[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg", 10, GREY, False)]])
        if logo:
            try:
                s.shapes.add_picture(logo, Inches(10.6), Inches(0.55), height=Inches(0.85))
            except Exception:
                pass
        self.page += 1
        return s

    def section(self, kicker, title, n, sub=""):
        s = self.slide()
        self.rect(s, 0, 0, self.SW, self.SH, WHITE)
        self.rect(s, 0, 0, Inches(0.28), self.SH, BLUE)
        self.rect(s, Inches(0.85), Inches(2.5), Inches(0.14), Inches(2.0), TEAL)
        self.txt(s, Inches(1.25), Inches(2.55), Inches(11), Inches(0.6), [[(kicker, 18, BLUE, True)]])
        self.txt(s, Inches(1.25), Inches(3.0), Inches(11.4), Inches(1.6), [[(title, 40, INK, True)]])
        if sub:
            self.txt(s, Inches(1.27), Inches(4.55), Inches(11), Inches(0.9), [[(sub, 16, GREY, False)]])
        if n:
            self.txt(s, Inches(10.0), Inches(0.7), Inches(2.8), Inches(1.6),
                     [[(n, 72, LINE, True)]], align=PP_ALIGN.RIGHT)
        self.footer(s)
        return s

    def content(self, title, items, kicker=None, size=20):
        s = self.head(self.slide(), title, kicker)
        self.bullets(s, Inches(0.85), Inches(1.95), Inches(11.6), Inches(4.9), items, size=size)
        self.footer(s); return s

    def two_col(self, title, left, right, kicker=None, lhead="", rhead="", lcolor=BLUE, rcolor=TEAL):
        s = self.head(self.slide(), title, kicker)
        self.rect(s, Inches(0.85), Inches(1.95), Inches(5.7), Inches(4.7), LIGHT)
        self.rect(s, Inches(6.95), Inches(1.95), Inches(5.55), Inches(4.7), LIGHT)
        self.rect(s, Inches(0.85), Inches(1.95), Inches(5.7), Inches(0.1), lcolor)
        self.rect(s, Inches(6.95), Inches(1.95), Inches(5.55), Inches(0.1), rcolor)
        if lhead:
            self.txt(s, Inches(1.1), Inches(2.2), Inches(5.2), Inches(0.4), [[(lhead, 16, lcolor, True)]])
        if rhead:
            self.txt(s, Inches(7.2), Inches(2.2), Inches(5.0), Inches(0.4), [[(rhead, 16, rcolor, True)]])
        self.bullets(s, Inches(1.1), Inches(2.75), Inches(5.2), Inches(3.75), left, size=15)
        self.bullets(s, Inches(7.2), Inches(2.75), Inches(5.05), Inches(3.75), right, size=15)
        self.footer(s); return s

    def cards3(self, title, cards, kicker):
        s = self.head(self.slide(), title, kicker)
        xs = [Inches(0.85), Inches(5.0), Inches(9.15)]
        for i, c in enumerate(cards[:3]):
            x = xs[i]; col = c[0]
            self.rect(s, x, Inches(1.95), Inches(3.65), Inches(4.7), LIGHT)
            self.rect(s, x, Inches(1.95), Inches(3.65), Inches(0.12), col)
            self.txt(s, x + Inches(0.25), Inches(2.2), Inches(3.2), Inches(0.6), [[(c[1], 19, col, True)]])
            self.bullets(s, x + Inches(0.25), Inches(2.95), Inches(3.2), Inches(3.4), c[2], size=14, gap=9)
        self.footer(s); return s

    def big_statement(self, line1, line2, kicker, color=BLUE):
        s = self.slide()
        self.rect(s, 0, 0, self.SW, self.SH, WHITE)
        self.rect(s, 0, 0, Inches(0.28), self.SH, color)
        self.txt(s, Inches(1.1), Inches(2.2), Inches(11), Inches(0.5), [[(kicker, 16, color, True)]])
        self.txt(s, Inches(1.1), Inches(2.8), Inches(11.3), Inches(2.4), [[(line1, 38, INK, True)]])
        if line2:
            self.txt(s, Inches(1.12), Inches(5.0), Inches(11), Inches(1.3), [[(line2, 20, GREY, False)]])
        self.footer(s); return s

    def tile_grid(self, title, items, kicker=None, cols=2, size=15, icons=None, accent=BLUE):
        s = self.head(self.slide(), title, kicker, kcolor=accent)
        n = len(items); rows = math.ceil(n / cols)
        X0 = Inches(0.85); Y0 = Inches(1.95); TOTW = Inches(11.63); AREAH = Inches(4.78)
        gx = Inches(0.3); gy = Inches(0.26)
        cw = int((TOTW - gx * (cols - 1)) / cols); ch = int((AREAH - gy * (rows - 1)) / rows)
        bd = Inches(0.58)
        for i, it in enumerate(items):
            r = i // cols; c = i % cols
            x = int(X0 + (cw + gx) * c); y = int(Y0 + (ch + gy) * r)
            col = PALETTE[i % len(PALETTE)]
            self.rect(s, x, y, cw, ch, LIGHT)
            self.rect(s, x, y, Inches(0.1), ch, col)
            self.oval(s, x + Inches(0.26), int(y + ch / 2 - bd / 2), bd, bd, col)
            ic = icons[i] if icons else str(i + 1)
            self.txt(s, x + Inches(0.26), int(y + ch / 2 - bd / 2), bd, bd,
                     [[(ic, 17, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            tx = x + Inches(1.02); tw = cw - Inches(1.26)
            if isinstance(it, tuple):
                self.txt(s, tx, int(y + Inches(0.12)), tw, int(ch - Inches(0.18)),
                         [[(it[0], size + 1, INK, True)], [(it[1], size - 3, GREY, False)]],
                         anchor=MSO_ANCHOR.MIDDLE, space=3)
            else:
                self.txt(s, tx, int(y + Inches(0.1)), tw, int(ch - Inches(0.16)),
                         [[(it, size, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
        self.footer(s); return s

    def flow_h(self, title, steps, kicker=None, color=BLUE, colors=None):
        s = self.head(self.slide(), title, kicker, kcolor=color)
        n = len(steps); X0 = Inches(0.85); TOTW = Inches(11.63); gap = Inches(0.34)
        cw = int((TOTW - gap * (n - 1)) / n)
        # Height follows the longest step text so short flows don't leave dead space.
        longest = max((len(str(t)) for t in steps), default=0)
        lines = max(1, -(-longest // max(1, int(cw / Inches(0.088)))))
        ch = int(min(Inches(3.15), Inches(1.95) + Inches(0.26) * lines))
        y = int(Inches(2.55) + (Inches(3.15) - ch) / 2)
        bd = Inches(0.8)
        for i, st in enumerate(steps):
            col = colors[i] if colors else color
            x = int(X0 + (cw + gap) * i)
            self.rect(s, x, y, cw, ch, LIGHT)
            self.rect(s, x, y, cw, Inches(0.1), col)
            self.oval(s, int(x + cw / 2 - bd / 2), int(y + Inches(0.42)), bd, bd, col)
            self.txt(s, int(x + cw / 2 - bd / 2), int(y + Inches(0.42)), bd, bd,
                     [[(str(i + 1), 28, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self.txt(s, x + Inches(0.14), int(y + Inches(1.42)), cw - Inches(0.28),
                     max(int(ch - Inches(1.54)), Inches(0.5)),
                     [[(st, 13, INK, False)]], align=PP_ALIGN.CENTER)
            if i < n - 1:
                self.txt(s, int(x + cw - Inches(0.04)), int(y + ch / 2 - Inches(0.3)),
                         int(gap + Inches(0.08)), Inches(0.6), [[("▶", 15, col, True)]],
                         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self.footer(s); return s

    def trainer_slide(self, kicker, name, role, rows, initials, accent=BLUE, photo=None):
        s = self.head(self.slide(), "About the Trainer", kicker, kcolor=accent)
        lx = Inches(0.85); lw = Inches(3.65)
        self.rect(s, lx, Inches(1.95), lw, Inches(4.7), LIGHT)
        self.rect(s, lx, Inches(1.95), lw, Inches(0.12), accent)
        bd = Inches(1.7); ax = int(lx + (lw - bd) / 2)
        placed = False
        if photo:
            try:
                s.shapes.add_picture(photo, ax, Inches(2.5), width=bd, height=bd); placed = True
            except Exception:
                placed = False
        if not placed:
            self.oval(s, ax, Inches(2.5), bd, bd, accent)
            self.txt(s, ax, Inches(2.5), bd, bd, [[(initials, 42, WHITE, True)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self.txt(s, lx + Inches(0.15), Inches(4.5), lw - Inches(0.3), Inches(0.6),
                 [[(name, 21, INK, True)]], align=PP_ALIGN.CENTER)
        self.txt(s, lx + Inches(0.15), Inches(5.15), lw - Inches(0.3), Inches(1.3),
                 [[(role, 13, GREY, False)]], align=PP_ALIGN.CENTER)
        rx = Inches(4.9); rw = Inches(7.6); ry = Inches(1.95); rh = Inches(4.7)
        n = len(rows); gy = Inches(0.18); th = int((rh - gy * (n - 1)) / n)
        for i, (label, val) in enumerate(rows):
            y = int(ry + (th + gy) * i); col = PALETTE[i % len(PALETTE)]
            self.rect(s, rx, y, rw, th, LIGHT)
            self.rect(s, rx, y, Inches(0.1), th, col)
            vruns = [(val, 13, INK, False)] if val else [("____________________________________________", 12, LINE, False)]
            self.txt(s, rx + Inches(0.3), y, rw - Inches(0.56), th,
                     [[(label.upper(), 10, col, True)], vruns], anchor=MSO_ANCHOR.MIDDLE, space=3)
        self.footer(s); return s

    def brk(self, kind, dur, color=AMBER):
        s = self.slide()
        self.rect(s, 0, 0, self.SW, self.SH, WHITE)
        self.rect(s, 0, 0, self.SW, Inches(0.22), color)
        self.rect(s, 0, Inches(7.28), self.SW, Inches(0.22), color)
        self.rect(s, Inches(5.4), Inches(2.35), Inches(2.53), Inches(0.1), color)
        self.txt(s, 0, Inches(2.75), self.SW, Inches(1.2), [[(kind, 48, INK, True)]], align=PP_ALIGN.CENTER)
        self.txt(s, 0, Inches(4.05), self.SW, Inches(0.8), [[(dur, 22, color, True)]], align=PP_ALIGN.CENTER)
        self.page += 1
        return s

    def image_slide(self, title, path, kicker=None, caption=None, accent=BLUE):
        """Full-bleed reference image with a caption strip."""
        s = self.head(self.slide(), title, kicker, kcolor=accent)
        top = Inches(1.9); avail_h = Inches(4.55) if caption else Inches(4.9)
        try:
            from PIL import Image
            with Image.open(path) as im:
                iw, ih = im.size
            ar = iw / ih
            max_w = Inches(11.2); max_h = avail_h
            w = max_w; h = int(w / ar)
            if h > max_h:
                h = max_h; w = int(h * ar)
            x = int((self.SW - w) / 2)
            s.shapes.add_picture(path, x, int(top), int(w), int(h))
            if caption:
                self.txt(s, Inches(0.85), int(top + h + Inches(0.12)), Inches(11.6), Inches(0.5),
                         [[(caption, 14, GREY, False)]], align=PP_ALIGN.CENTER)
        except Exception:
            self.txt(s, Inches(0.85), Inches(3.0), Inches(11.6), Inches(1.0),
                     [[(caption or "", 16, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(s); return s

    # ================= Lean Six Sigma concept diagrams =================
    def dmaic_wheel(self, title, phases, kicker=None):
        """DMAIC as five connected chevrons with per-phase detail beneath."""
        s = self.head(self.slide(), title, kicker, kcolor=BLUE)
        n = len(phases); X0 = Inches(0.85); TOTW = Inches(11.63)
        cw = int(TOTW / n); y = Inches(2.05); ch = Inches(0.95)
        for i, (letter, name, detail) in enumerate(phases):
            col = DMAIC_COLORS[i % len(DMAIC_COLORS)]
            x = int(X0 + cw * i)
            self.chevron(s, x, y, int(cw * 1.02), ch, col)
            self.txt(s, x + Inches(0.16), y, int(cw * 0.86), ch,
                     [[(letter, 26, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            py = int(y + ch + Inches(0.34))
            self.rect(s, x + Inches(0.06), py, int(cw - Inches(0.16)), Inches(3.35), LIGHT)
            self.rect(s, x + Inches(0.06), py, int(cw - Inches(0.16)), Inches(0.09), col)
            self.txt(s, x + Inches(0.22), int(py + Inches(0.22)), int(cw - Inches(0.48)), Inches(0.55),
                     [[(name, 15, col, True)]])
            self.bullets(s, x + Inches(0.22), int(py + Inches(0.82)), int(cw - Inches(0.46)),
                         Inches(2.4), detail, size=11, gap=6)
        self.footer(s); return s

    def sipoc_diagram(self, title, cols, kicker=None):
        """SIPOC as five colour-coded columns feeding left to right."""
        s = self.head(self.slide(), title, kicker, kcolor=TEAL)
        names = ["Suppliers", "Inputs", "Process", "Outputs", "Customers"]
        letters = ["S", "I", "P", "O", "C"]
        cols_c = [BLUE, TEAL, VIOLET, AMBER, RED]
        n = 5; X0 = Inches(0.85); TOTW = Inches(11.63); gap = Inches(0.16)
        cw = int((TOTW - gap * (n - 1)) / n); y = Inches(2.0); ch = Inches(3.55)
        for i in range(n):
            x = int(X0 + (cw + gap) * i); col = cols_c[i]
            self.rect(s, x, y, cw, ch, LIGHT)
            self.rect(s, x, y, cw, Inches(0.62), col)
            self.txt(s, x, y, cw, Inches(0.62),
                     [[(f"{letters[i]} — {names[i]}", 14, WHITE, True)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self.bullets(s, x + Inches(0.16), int(y + Inches(0.8)), cw - Inches(0.32),
                         int(ch - Inches(0.92)), cols[i], size=12, gap=7)
            if i < n - 1:
                self.txt(s, int(x + cw), int(y + ch / 2 - Inches(0.25)), int(gap), Inches(0.5),
                         [[("▶", 12, GREY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self.txt(s, Inches(0.85), Inches(6.05), Inches(11.63), Inches(0.5),
                 [[("Read left to right: suppliers provide the inputs, the process transforms them, "
                    "and the outputs go to the customers.", 13, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(s); return s

    def fishbone(self, title, effect, bones, kicker=None):
        """Ishikawa / fishbone: spine to the effect box with angled category bones.

        Bones alternate above/below the spine. Category labels sit adjacent to
        the spine and cause text is placed OUTWARD from the label so it never
        collides with the slide title or the footer.
        """
        s = self.head(self.slide(), title, kicker, kcolor=VIOLET)
        cy = Inches(4.32)
        # Spine spans from the left margin to the effect box, sized to the number
        # of bones so it never trails past the last one.
        n_top = len(bones[0::2])
        last_bone_x = Inches(2.55) + Inches(2.85) * (n_top - 1) + Inches(0.7)
        head_x = max(Inches(9.4), int(last_bone_x + Inches(0.5)))
        x0 = Inches(0.85)
        self.line_seg(s, x0, cy, head_x, cy, INK, width=2.5)
        # effect box (the head of the fish)
        self.rect(s, head_x, Inches(3.62), Inches(2.35), Inches(1.4), RED)
        self.txt(s, head_x, Inches(3.62), Inches(2.35), Inches(1.4),
                 [[(effect, 13, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        top = bones[0::2]; bot = bones[1::2]
        LBL_W = Inches(1.75); LBL_H = Inches(0.4)
        for grp, above in ((top, True), (bot, False)):
            n = len(grp)
            for j, (cat, causes) in enumerate(grp):
                # spread bones evenly along the usable spine length
                bx = int(Inches(2.55) + Inches(2.85) * j)
                col = PALETTE[(j * 2 + (0 if above else 1)) % len(PALETTE)]
                if above:
                    lbl_y = int(cy - Inches(1.28))          # label near the spine
                    txt_y = int(cy - Inches(2.28))          # causes further out (up)
                else:
                    lbl_y = int(cy + Inches(0.88))
                    txt_y = int(cy + Inches(1.34))
                # angled bone from the label into the spine
                bone_from_y = int(lbl_y + LBL_H) if above else lbl_y
                self.line_seg(s, bx, bone_from_y, int(bx + Inches(0.7)), cy, col, width=2)
                self.rect(s, int(bx - LBL_W / 2), lbl_y, LBL_W, LBL_H, col)
                self.txt(s, int(bx - LBL_W / 2), lbl_y, LBL_W, LBL_H,
                         [[(cat, 12, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                self.bullets(s, int(bx - Inches(1.05)), txt_y, Inches(2.35), Inches(0.95),
                             causes, size=10, gap=1)
        self.footer(s); return s

    def pareto_chart(self, title, data, kicker=None, note=None):
        """Pareto: descending bars + cumulative % line, drawn to scale."""
        s = self.head(self.slide(), title, kicker, kcolor=AMBER)
        data = sorted(data, key=lambda d: -d[1])
        total = sum(v for _, v in data) or 1
        # plot area leaves room for value labels above and 2-line category labels below
        PX = Inches(1.35); PY = Inches(2.35); PW = Inches(9.9); PH = Inches(3.1)
        self.line_seg(s, PX, PY + PH, PX + PW, PY + PH, GREY, width=1.5)
        self.line_seg(s, PX, PY, PX, PY + PH, GREY, width=1.5)
        n = len(data); slot = PW / n; bw = slot * 0.6
        mx = max(v for _, v in data)
        cum = 0; pts = []
        for i, (label, v) in enumerate(data):
            h = int(PH * (v / mx) * 0.9)
            x = int(PX + slot * i + (slot - bw) / 2)
            y = int(PY + PH - h)
            self.rect(s, x, y, int(bw), h, BLUE)
            self.txt(s, int(PX + slot * i), int(y - Inches(0.32)), int(slot), Inches(0.3),
                     [[(str(v), 11, INK, True)]], align=PP_ALIGN.CENTER)
            self.txt(s, int(PX + slot * i - Inches(0.1)), int(PY + PH + Inches(0.12)),
                     int(slot + Inches(0.2)), Inches(0.72),
                     [[(label, 10, GREY, False)]], align=PP_ALIGN.CENTER)
            cum += v
            pts.append((int(PX + slot * i + slot / 2), int(PY + PH - PH * (cum / total) * 0.9), cum / total))
        for i in range(len(pts) - 1):
            self.line_seg(s, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], RED, width=2)
        for i, (px, py, frac) in enumerate(pts):
            self.oval(s, px - Inches(0.07), py - Inches(0.07), Inches(0.14), Inches(0.14), RED)
            # Early points sit low on the curve, right where the tall bars are. Place
            # their labels to the LEFT of the marker (in the gap between bars) instead
            # of above it, so they never collide with a bar's value label.
            bar_top = int(PY + PH - PH * (data[i][1] / mx) * 0.9)
            crowded = py > bar_top - Inches(0.5)
            if crowded:
                lx = px - Inches(1.05); ly = py - Inches(0.15)
            else:
                lx = px - Inches(0.45); ly = py - Inches(0.48)
            self.txt(s, lx, ly, Inches(0.9), Inches(0.3),
                     [[(f"{frac*100:.0f}%", 10, RED, True)]], align=PP_ALIGN.CENTER)
        # 80% reference line
        y80 = int(PY + PH - PH * 0.8 * 0.9)
        self.line_seg(s, PX, y80, PX + PW, y80, GREY, width=1, dash=True)
        self.txt(s, int(PX + PW + Inches(0.06)), y80 - Inches(0.16), Inches(0.9), Inches(0.3),
                 [[("80%", 10, GREY, True)]])
        self.txt(s, Inches(0.85), Inches(6.45), Inches(11.6), Inches(0.5),
                 [[(note or "The vital few: the leftmost categories drive most of the problem — fix those first.",
                    13, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(s); return s

    def run_chart(self, title, series, kicker=None, note=None, median_label="Median"):
        """Run chart: points over time with a median reference line."""
        s = self.head(self.slide(), title, kicker, kcolor=TEAL)
        PX = Inches(1.3); PY = Inches(2.15); PW = Inches(10.3); PH = Inches(3.6)
        self.line_seg(s, PX, PY + PH, PX + PW, PY + PH, GREY, width=1.5)
        self.line_seg(s, PX, PY, PX, PY + PH, GREY, width=1.5)
        vals = [v for _, v in series]
        lo, hi = min(vals), max(vals)
        rng = (hi - lo) or 1
        lo_p = lo - rng * 0.15; hi_p = hi + rng * 0.15
        span = hi_p - lo_p

        def ypix(v):
            return int(PY + PH - PH * ((v - lo_p) / span))
        n = len(series); step = PW / max(n - 1, 1)
        pts = [(int(PX + step * i), ypix(v)) for i, (_, v) in enumerate(series)]
        sv = sorted(vals); m = sv[len(sv) // 2] if len(sv) % 2 else (sv[len(sv) // 2 - 1] + sv[len(sv) // 2]) / 2
        my = ypix(m)
        self.line_seg(s, PX, my, PX + PW, my, AMBER, width=1.6, dash=True)
        self.txt(s, int(PX + PW + Inches(0.04)), my - Inches(0.18), Inches(1.4), Inches(0.36),
                 [[(f"{median_label} {m:g}", 10, AMBER, True)]])
        for i in range(len(pts) - 1):
            self.line_seg(s, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], BLUE, width=2)
        for i, (px, py) in enumerate(pts):
            self.oval(s, px - Inches(0.075), py - Inches(0.075), Inches(0.15), Inches(0.15), BLUE)
            if n <= 16:
                self.txt(s, px - Inches(0.4), int(PY + PH + Inches(0.1)), Inches(0.8), Inches(0.4),
                         [[(str(series[i][0]), 9, GREY, False)]], align=PP_ALIGN.CENTER)
        self.txt(s, Inches(0.85), Inches(6.2), Inches(11.6), Inches(0.5),
                 [[(note or "Look for trend, shift, cluster and oscillation — never react to a single point.",
                    13, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(s); return s

    def normal_curve(self, title, kicker=None, note=None, sigma_labels=True):
        """Bell curve with ±sigma bands and spec limits — the heart of 'six sigma'."""
        s = self.head(self.slide(), title, kicker, kcolor=VIOLET)
        CX = Inches(6.65); BASE = Inches(5.65); AMP = Inches(3.0); SD = Inches(1.35)
        self.line_seg(s, Inches(1.2), BASE, Inches(12.1), BASE, GREY, width=1.5)
        # curve as a polyline of small segments
        prev = None
        steps = 120
        for i in range(steps + 1):
            z = -4.0 + 8.0 * i / steps
            x = int(CX + SD * z)
            y = int(BASE - AMP * math.exp(-0.5 * z * z))
            if prev:
                self.line_seg(s, prev[0], prev[1], x, y, VIOLET, width=2)
            prev = (x, y)
        # sigma gridlines
        for z in (-3, -2, -1, 0, 1, 2, 3):
            x = int(CX + SD * z)
            h = AMP * math.exp(-0.5 * z * z)
            self.line_seg(s, x, BASE, x, int(BASE - h), LINE, width=1)
            if sigma_labels:
                lab = "μ" if z == 0 else f"{z:+d}σ"
                self.txt(s, x - Inches(0.35), int(BASE + Inches(0.1)), Inches(0.7), Inches(0.3),
                         [[(lab, 11, GREY, True)]], align=PP_ALIGN.CENTER)
        # spec limits at ±3σ
        for z, lab, col in ((-3, "LSL", RED), (3, "USL", RED)):
            x = int(CX + SD * z)
            self.line_seg(s, x, int(BASE - AMP * 1.02), x, BASE, col, width=1.8, dash=True)
            self.txt(s, x - Inches(0.5), int(BASE - AMP * 1.18), Inches(1.0), Inches(0.32),
                     [[(lab, 12, col, True)]], align=PP_ALIGN.CENTER)
        self.txt(s, Inches(0.85), Inches(6.25), Inches(11.6), Inches(0.6),
                 [[(note or "The wider the spec limits sit from the mean in sigma units, the fewer defects the process produces.",
                    14, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(s); return s

    def waste_wheel(self, title, wastes, kicker=None, note=None):
        """The eight wastes as a labelled ring around a central acronym."""
        s = self.head(self.slide(), title, kicker, kcolor=RED)
        CX = Inches(6.66); CY = Inches(4.0); R = Inches(1.95); bd = Inches(1.34)
        self.oval(s, CX - Inches(1.02), CY - Inches(1.02), Inches(2.04), Inches(2.04), LIGHT)
        self.txt(s, CX - Inches(1.02), CY - Inches(1.02), Inches(2.04), Inches(2.04),
                 [[("DOWNTIME", 17, RED, True)], [("8 wastes", 12, GREY, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=2)
        n = len(wastes)
        for i, (letter, name) in enumerate(wastes):
            ang = -math.pi / 2 + 2 * math.pi * i / n
            x = int(CX + R * math.cos(ang) * 1.62 - bd / 2)
            y = int(CY + R * math.sin(ang) * 0.92 - bd / 2)
            col = PALETTE[i % len(PALETTE)]
            self.oval(s, x, y, bd, bd, col)
            self.txt(s, x, y, bd, bd, [[(letter, 20, WHITE, True)], [(name, 8, WHITE, False)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=1)
        self.txt(s, Inches(0.85), Inches(6.4), Inches(11.6), Inches(0.5),
                 [[(note or "Waste is any effort the customer would not willingly pay for.", 14, GREY, False)]],
                 align=PP_ALIGN.CENTER)
        self.footer(s); return s

    def vs_diagram(self, title, left, right, combined, kicker=None):
        """A vs B → combined. Used for Lean vs Six Sigma → Lean Six Sigma."""
        s = self.head(self.slide(), title, kicker, kcolor=BLUE)
        bw = Inches(4.35); bh = Inches(2.5); y = Inches(2.0)
        self.rect(s, Inches(0.85), y, bw, bh, LIGHT)
        self.rect(s, Inches(0.85), y, bw, Inches(0.1), BLUE)
        self.txt(s, Inches(1.1), y + Inches(0.28), bw - Inches(0.5), Inches(0.5), [[(left[0], 22, BLUE, True)]])
        self.bullets(s, Inches(1.1), y + Inches(0.95), bw - Inches(0.5), Inches(1.4), left[1], size=13, gap=7)
        self.rect(s, Inches(8.1), y, bw, bh, LIGHT)
        self.rect(s, Inches(8.1), y, bw, Inches(0.1), TEAL)
        self.txt(s, Inches(8.35), y + Inches(0.28), bw - Inches(0.5), Inches(0.5), [[(right[0], 22, TEAL, True)]])
        self.bullets(s, Inches(8.35), y + Inches(0.95), bw - Inches(0.5), Inches(1.4), right[1], size=13, gap=7)
        self.txt(s, Inches(5.6), y + Inches(0.85), Inches(2.1), Inches(0.8),
                 [[("+", 44, GREY, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self.txt(s, Inches(0.85), Inches(4.72), Inches(11.6), Inches(0.4),
                 [[("▼", 16, GREY, True)]], align=PP_ALIGN.CENTER)
        self.rect(s, Inches(2.4), Inches(5.15), Inches(8.5), Inches(1.42), VIOLET)
        self.txt(s, Inches(2.6), Inches(5.15), Inches(8.1), Inches(1.42),
                 [[(combined[0], 22, WHITE, True)], [(combined[1], 13, WHITE, False)]],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space=4)
        self.footer(s); return s

    def matrix2x2(self, title, xlab, ylab, quads, kicker=None, accent=BLUE):
        """2x2 matrix — impact/effort, Kano-style positioning, etc."""
        s = self.head(self.slide(), title, kicker, kcolor=accent)
        X0 = Inches(2.5); Y0 = Inches(2.0); W = Inches(8.4); H = Inches(4.1)
        hw = int(W / 2); hh = int(H / 2)
        cols = [TEAL, BLUE, AMBER, GREY]
        pos = [(0, 0), (1, 0), (0, 1), (1, 1)]
        for i, q in enumerate(quads):
            cx, ry = pos[i]
            x = int(X0 + hw * cx); y = int(Y0 + hh * ry)
            self.rect(s, x, y, hw, hh, LIGHT, line=LINE)
            self.rect(s, x, y, Inches(0.08), hh, cols[i])
            self.txt(s, x + Inches(0.28), y + Inches(0.24), hw - Inches(0.56), Inches(0.5),
                     [[(q[0], 16, cols[i], True)]])
            self.txt(s, x + Inches(0.28), y + Inches(0.78), hw - Inches(0.56), hh - Inches(1.0),
                     [[(q[1], 12, GREY, False)]])
        self.txt(s, X0, int(Y0 + H + Inches(0.12)), W, Inches(0.4),
                 [[(xlab, 13, INK, True)]], align=PP_ALIGN.CENTER)
        tb = self.txt(s, Inches(0.6), int(Y0 + H / 2 - Inches(0.2)), Inches(1.8), Inches(0.4),
                      [[(ylab, 13, INK, True)]], align=PP_ALIGN.CENTER)
        self.footer(s); return s

    def ladder(self, title, rungs, kicker=None, accent=BLUE, note=None):
        """Ascending stacked bars — belt pathway, sigma levels, maturity."""
        s = self.head(self.slide(), title, kicker, kcolor=accent)
        n = len(rungs); X0 = Inches(1.0); TOTW = Inches(11.3); gap = Inches(0.22)
        cw = int((TOTW - gap * (n - 1)) / n)
        base = Inches(6.05); maxh = Inches(3.85)
        for i, (label, sub) in enumerate(rungs):
            h = int(maxh * (0.34 + 0.66 * (i + 1) / n))
            x = int(X0 + (cw + gap) * i); y = int(base - h)
            col = PALETTE[i % len(PALETTE)]
            self.rect(s, x, y, cw, h, LIGHT)
            self.rect(s, x, y, cw, Inches(0.09), col)
            self.txt(s, x + Inches(0.12), int(y + Inches(0.24)), cw - Inches(0.24), Inches(0.6),
                     [[(label, 15, col, True)]], align=PP_ALIGN.CENTER)
            self.txt(s, x + Inches(0.14), int(y + Inches(0.88)), cw - Inches(0.28), int(h - Inches(1.0)),
                     [[(sub, 11, GREY, False)]], align=PP_ALIGN.CENTER)
        self.line_seg(s, X0, base, Inches(12.3), base, GREY, width=1.5)
        if note:
            self.txt(s, Inches(0.85), Inches(6.3), Inches(11.6), Inches(0.5),
                     [[(note, 13, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(s); return s

    def formula_card(self, title, items, kicker=None, accent=VIOLET, note=None):
        """Formula panels: name, expression, worked example."""
        s = self.head(self.slide(), title, kicker, kcolor=accent)
        n = len(items); Y0 = Inches(1.95); H = Inches(4.6); gy = Inches(0.2)
        rh = int((H - gy * (n - 1)) / n)
        for i, (name, expr, ex) in enumerate(items):
            y = int(Y0 + (rh + gy) * i); col = PALETTE[i % len(PALETTE)]
            self.rect(s, Inches(0.85), y, Inches(11.63), rh, LIGHT)
            self.rect(s, Inches(0.85), y, Inches(0.1), rh, col)
            self.txt(s, Inches(1.15), y, Inches(2.9), rh, [[(name, 15, col, True)]], anchor=MSO_ANCHOR.MIDDLE)
            self.rect(s, Inches(4.1), int(y + Inches(0.12)), Inches(4.3), int(rh - Inches(0.24)), WHITE, line=LINE)
            self.txt(s, Inches(4.2), y, Inches(4.1), rh, [[(expr, 14, INK, True)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            self.txt(s, Inches(8.65), y, Inches(3.6), rh, [[(ex, 12, GREY, False)]], anchor=MSO_ANCHOR.MIDDLE)
        if note:
            self.txt(s, Inches(0.85), Inches(6.65), Inches(11.6), Inches(0.4),
                     [[(note, 12, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(s); return s

    def timeline(self, title, events, kicker=None, accent=AMBER):
        """Horizontal history timeline."""
        s = self.head(self.slide(), title, kicker, kcolor=accent)
        # Inset the first/last nodes by half a card so their cards stay on-slide.
        CARD_W = Inches(2.1); half = int(CARD_W / 2)
        X0 = Inches(0.85) + half; W = Inches(11.63) - CARD_W
        y = Inches(4.1)
        self.line_seg(s, Inches(0.85), y, Inches(0.85) + Inches(11.63), y, LINE, width=3)
        n = len(events); step = W / max(n - 1, 1)
        for i, (yr, ev) in enumerate(events):
            x = int(X0 + step * i); col = PALETTE[i % len(PALETTE)]
            self.oval(s, x - Inches(0.16), y - Inches(0.16), Inches(0.32), Inches(0.32), col)
            above = (i % 2 == 0)
            by = int(y - Inches(1.62)) if above else int(y + Inches(0.34))
            self.rect(s, x - Inches(1.05), by, Inches(2.1), Inches(1.28), LIGHT)
            self.rect(s, x - Inches(1.05), by, Inches(2.1), Inches(0.08), col)
            self.txt(s, x - Inches(0.95), int(by + Inches(0.18)), Inches(1.9), Inches(0.36),
                     [[(yr, 15, col, True)]], align=PP_ALIGN.CENTER)
            self.txt(s, x - Inches(0.98), int(by + Inches(0.58)), Inches(1.96), Inches(0.64),
                     [[(ev, 10, GREY, False)]], align=PP_ALIGN.CENTER)
        self.footer(s); return s

    def compare_panels(self, title, panels, kicker=None, accent=BLUE):
        """N side-by-side labelled panels (2-4) for concept comparisons."""
        s = self.head(self.slide(), title, kicker, kcolor=accent)
        n = len(panels); X0 = Inches(0.85); TOTW = Inches(11.63); gap = Inches(0.28)
        cw = int((TOTW - gap * (n - 1)) / n); y = Inches(1.95); ch = Inches(4.7)
        for i, (hd, sub, items) in enumerate(panels):
            x = int(X0 + (cw + gap) * i); col = PALETTE[i % len(PALETTE)]
            self.rect(s, x, y, cw, ch, LIGHT)
            self.rect(s, x, y, cw, Inches(0.1), col)
            self.txt(s, x + Inches(0.22), y + Inches(0.28), cw - Inches(0.44), Inches(0.5),
                     [[(hd, 18, col, True)]])
            self.txt(s, x + Inches(0.22), y + Inches(0.82), cw - Inches(0.44), Inches(0.5),
                     [[(sub, 12, GREY, False)]])
            self.bullets(s, x + Inches(0.22), y + Inches(1.42), cw - Inches(0.44), Inches(3.0),
                         items, size=12, gap=7)
        self.footer(s); return s

    # ---------------- activity slides ----------------
    def activity_overview(self, tag, title, desc, build, tools, kicker, elective=False):
        s = self.head(self.slide(), title, kicker, kcolor=TEAL)
        self.rect(s, Inches(0.85), Inches(1.85), Inches(1.7), Inches(0.5), TEAL)
        self.txt(s, Inches(0.85), Inches(1.9), Inches(1.7), Inches(0.4),
                 [[(tag, 16, WHITE, True)]], align=PP_ALIGN.CENTER)
        if elective:
            self.rect(s, Inches(2.7), Inches(1.85), Inches(1.6), Inches(0.5), AMBER)
            self.txt(s, Inches(2.7), Inches(1.9), Inches(1.6), Inches(0.4),
                     [[("ELECTIVE", 13, WHITE, True)]], align=PP_ALIGN.CENTER)
        self.txt(s, Inches(0.85), Inches(2.55), Inches(11.7), Inches(1.6), [[(desc, 20, INK, False)]])
        self.rect(s, Inches(0.85), Inches(4.3), Inches(11.7), Inches(2.0), LIGHT)
        self.txt(s, Inches(1.1), Inches(4.5), Inches(11), Inches(0.4), [[("You'll build", 14, BLUE, True)]])
        self.txt(s, Inches(1.1), Inches(4.9), Inches(11), Inches(0.6), [[(build, 18, INK, True)]])
        self.txt(s, Inches(1.1), Inches(5.6), Inches(11.2), Inches(0.6),
                 [[("Tools & techniques:  ", 13, GREY, True), (tools, 13, GREY, False)]])
        self.footer(s); return s

    def step_slide(self, kicker, act_title, n, total, text, cmd=""):
        s = self.head(self.slide(), act_title, kicker, TEAL)
        self.oval(s, Inches(0.85), Inches(2.5), Inches(1.4), Inches(1.4), TEAL)
        self.txt(s, Inches(0.85), Inches(2.5), Inches(1.4), Inches(1.4),
                 [[(str(n), 36, WHITE, True)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        self.txt(s, Inches(0.95), Inches(1.95), Inches(11), Inches(0.4),
                 [[(f"STEP {n} OF {total}", 13, GREY, True)]])
        self.txt(s, Inches(2.55), Inches(2.4), Inches(10.1), Inches(1.6),
                 [[(text, 22, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
        if cmd:
            self.rect(s, Inches(2.55), Inches(4.3), Inches(10.1), Inches(0.95), RGBColor(0x0B, 0x12, 0x20))
            self.txt(s, Inches(2.8), Inches(4.3), Inches(9.7), Inches(0.95),
                     [[(cmd, 13, RGBColor(0x9C, 0xDC, 0xFE), False)]], anchor=MSO_ANCHOR.MIDDLE)
        self.footer(s); return s

    def steps_group(self, kicker, act_title, items, first, total):
        """Several lab steps on one slide — numbered rows with optional command box.

        `items` is a list of (instruction, cmd) tuples; `first` is the 1-based
        number of items[0] within the lab, and `total` the lab's step count.
        Keeps every step's wording while holding the deck to a sane length.
        """
        s = self.head(self.slide(), act_title, kicker, TEAL)
        last = first + len(items) - 1
        self.txt(s, Inches(0.9), Inches(1.86), Inches(11), Inches(0.36),
                 [[(f"STEPS {first}-{last} OF {total}", 13, GREY, True)]])
        n = len(items)
        Y0 = Inches(2.32); AREA = Inches(4.35); gy = Inches(0.16)
        rh = int((AREA - gy * (n - 1)) / max(n, 1))
        bd = Inches(0.5)
        for i, it in enumerate(items):
            instr, cmd = (it if isinstance(it, (tuple, list)) else (it, ""))
            y = int(Y0 + (rh + gy) * i)
            self.rect(s, Inches(0.85), y, Inches(11.63), rh, LIGHT)
            self.rect(s, Inches(0.85), y, Inches(0.09), rh, TEAL)
            self.oval(s, Inches(1.12), int(y + rh / 2 - bd / 2), bd, bd, TEAL)
            self.txt(s, Inches(1.12), int(y + rh / 2 - bd / 2), bd, bd,
                     [[(str(first + i), 15, WHITE, True)]],
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            tx = Inches(1.86)
            tw = Inches(10.4) if not cmd else Inches(6.05)
            self.txt(s, tx, int(y + Inches(0.06)), tw, int(rh - Inches(0.12)),
                     [[(instr, 13, INK, False)]], anchor=MSO_ANCHOR.MIDDLE)
            if cmd:
                cx = Inches(8.05)
                self.rect(s, cx, int(y + Inches(0.12)), Inches(4.3), int(rh - Inches(0.24)),
                          RGBColor(0x0B, 0x12, 0x20))
                self.txt(s, cx + Inches(0.14), int(y + Inches(0.12)), Inches(4.02),
                         int(rh - Inches(0.24)),
                         [[(cmd, 10, RGBColor(0x9C, 0xDC, 0xFE), False)]],
                         anchor=MSO_ANCHOR.MIDDLE)
        self.footer(s); return s

    def test_slide(self, act_title, text, kicker):
        s = self.head(self.slide(), act_title, kicker, TEAL)
        self.rect(s, Inches(0.85), Inches(2.3), Inches(11.7), Inches(2.6), RGBColor(0xE8, 0xF7, 0xEE))
        self.txt(s, Inches(1.2), Inches(2.6), Inches(11), Inches(0.5),
                 [[("✅  Check your work", 20, RGBColor(0x12, 0x7A, 0x3E), True)]])
        self.txt(s, Inches(1.2), Inches(3.3), Inches(11), Inches(1.4), [[(text, 18, INK, False)]])
        self.footer(s); return s

    def checkpoint(self, act_title, questions, kicker):
        s = self.head(self.slide(), act_title, kicker, kcolor=VIOLET)
        self.bullets(s, Inches(0.95), Inches(2.2), Inches(11.4), Inches(4.2),
                     questions, size=19, gap=18)
        self.footer(s); return s

    # ---------------- transitions ----------------
    def apply_transitions(self, kind="fade", dur_ms=700):
        """Apply a slide transition to every slide.

        python-pptx has no transition API, so the PowerPoint transition element
        is injected directly into each slide's XML. Uses the p14 (PowerPoint
        2010+) namespace with mc:AlternateContent so older renderers degrade
        gracefully instead of erroring.
        """
        for sl in self.prs.slides:
            _set_transition(sl, kind, dur_ms)


# ---------------- XML helpers ----------------
from pptx.oxml import parse_xml

P14 = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
MC = 'http://schemas.openxmlformats.org/markup-compatibility/2006'
A = 'http://schemas.openxmlformats.org/drawingml/2006/main'
P = 'http://schemas.openxmlformats.org/presentationml/2006/main'


def _el(tag, attrs=None):
    from lxml import etree
    ns = {'a': A}
    prefix, local = tag.split(':')
    e = etree.SubElement(etree.Element('root', nsmap=ns), f'{{{A}}}{local}')
    if attrs:
        for k, v in attrs.items():
            e.set(k, v)
    return e


_TRANSITIONS = {
    # name -> (element name, extra attributes)
    "fade":  ("p:fade", {}),
    "push":  ("p:push", {"dir": "u"}),
    "wipe":  ("p:wipe", {"dir": "d"}),
    "morph": ("p14:morph", {"option": "byObject"}),
}


def _set_transition(slide, kind="fade", dur_ms=700):
    """Insert <mc:AlternateContent><p:transition> into the slide XML."""
    sld = slide._element
    # remove any existing transition
    for old in sld.findall(f'{{{P}}}transition') + sld.findall(f'{{{MC}}}AlternateContent'):
        sld.remove(old)

    elname, attrs = _TRANSITIONS.get(kind, _TRANSITIONS["fade"])
    attr_str = " ".join(f'{k}="{v}"' for k, v in attrs.items())
    inner = f'<{elname} {attr_str}/>' if attr_str else f'<{elname}/>'

    xml = (
        f'<mc:AlternateContent xmlns:mc="{MC}" xmlns:p14="{P14}" '
        f'xmlns:p="{P}" xmlns:a="{A}">'
        f'<mc:Choice xmlns:p14="{P14}" Requires="p14">'
        f'<p:transition spd="med" p14:dur="{dur_ms}">{inner}</p:transition>'
        f'</mc:Choice>'
        f'<mc:Fallback>'
        f'<p:transition spd="med"><p:fade/></p:transition>'
        f'</mc:Fallback>'
        f'</mc:AlternateContent>'
    )
    frag = parse_xml(xml)
    # transition must come after the cSld/clrMapOvr elements
    sld.append(frag)
    return frag
